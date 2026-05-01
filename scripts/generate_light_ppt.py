from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


FONT_HEAD = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_CODE = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SOURCE_TXT = Path("新建 文本文档.txt")
OUT_DIR = Path("artifacts/ppt")
OUT_PATH = OUT_DIR / "odcr_preprocess_bright_report.pptx"
SINGLE_SLIDE_OUT_PATH = OUT_DIR / "odcr_overall_innovation_chain_single_slide.pptx"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.replace("#", ""))


PALETTE = {
    "bg": rgb("#F7FAFF"),
    "surface": rgb("#FFFFFF"),
    "surface_alt": rgb("#F1F7FF"),
    "border": rgb("#D7E5F5"),
    "text": rgb("#1E2A46"),
    "muted": rgb("#5D6B8A"),
    "blue": rgb("#7BC6FF"),
    "blue_soft": rgb("#EAF6FF"),
    "mint": rgb("#A7E6D5"),
    "mint_soft": rgb("#EEFBF7"),
    "gold": rgb("#FFE29A"),
    "gold_soft": rgb("#FFF9E6"),
    "coral": rgb("#FFC7B0"),
    "coral_soft": rgb("#FFF2EC"),
    "lavender": rgb("#D7CCFF"),
    "lavender_soft": rgb("#F4F1FF"),
}

ACCENTS = ["blue", "mint", "gold", "coral", "lavender"]


@dataclass
class BoxStyle:
    fill: str = "surface"
    accent: str = "blue"
    title_size: int = 20
    body_size: int = 16


def add_full_background(slide, accent_idx: int = 0) -> None:
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE["bg"]
    bg.line.fill.background()

    accents = [ACCENTS[accent_idx % len(ACCENTS)], ACCENTS[(accent_idx + 2) % len(ACCENTS)], ACCENTS[(accent_idx + 3) % len(ACCENTS)]]

    top_right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.5), Inches(-0.65), Inches(3.1), Inches(3.1))
    top_right.fill.solid()
    top_right.fill.fore_color.rgb = PALETTE[f"{accents[0]}_soft"]
    top_right.line.fill.background()

    bottom_left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.7), Inches(5.9), Inches(2.6), Inches(2.6))
    bottom_left.fill.solid()
    bottom_left.fill.fore_color.rgb = PALETTE[f"{accents[1]}_soft"]
    bottom_left.line.fill.background()

    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.65), Inches(5.95), Inches(2.9), Inches(0.5))
    band.fill.solid()
    band.fill.fore_color.rgb = PALETTE[f"{accents[2]}_soft"]
    band.line.color.rgb = PALETTE["border"]


def set_shape_fill(shape, fill_key: str, line_key: str = "border") -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALETTE[fill_key]
    shape.line.color.rgb = PALETTE[line_key]
    shape.line.width = Pt(1)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    paragraphs: Sequence[dict],
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Pt(4)
    frame.margin_right = Pt(4)
    frame.margin_top = Pt(2)
    frame.margin_bottom = Pt(2)
    for idx, spec in enumerate(paragraphs):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = spec["text"]
        p.font.name = spec.get("font", FONT_BODY)
        p.font.size = Pt(spec.get("size", 16))
        p.font.bold = spec.get("bold", False)
        p.font.italic = spec.get("italic", False)
        p.font.color.rgb = spec.get("color", PALETTE["text"])
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(spec.get("space_after", 4))
        p.space_before = Pt(spec.get("space_before", 0))
        if spec.get("level") is not None:
            p.level = spec["level"]
        if spec.get("bullet", False):
            p.text = f"• {p.text}"
    return box


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    body_lines: Iterable[str],
    style: BoxStyle,
    title_color=None,
):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_shape_fill(card, style.fill)

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.09), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = PALETTE[style.accent]
    stripe.line.fill.background()

    paragraphs = [
        {"text": title, "size": style.title_size, "bold": True, "color": title_color or PALETTE["text"], "space_after": 8},
    ]
    paragraphs.extend({"text": line, "size": style.body_size, "color": PALETTE["muted"], "space_after": 3} for line in body_lines)
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.28), height - Inches(0.2), paragraphs)
    return card


def add_pill(slide, left, top, width, height, text: str, fill_key: str, text_color=None, font_size: int = 14):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_shape_fill(pill, fill_key, line_key=fill_key)
    add_text(
        slide,
        left + Inches(0.03),
        top + Inches(0.01),
        width - Inches(0.06),
        height - Inches(0.03),
        [{"text": text, "size": font_size, "bold": True, "align": PP_ALIGN.CENTER, "color": text_color or PALETTE["text"], "space_after": 0}],
        valign=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    return pill


def add_slide_title(slide, title: str, subtitle: str | None, accent_idx: int, slide_no: int, total: int) -> None:
    add_full_background(slide, accent_idx)
    add_text(
        slide,
        Inches(0.7),
        Inches(0.35),
        Inches(10.5),
        Inches(0.55),
        [{"text": title, "size": 28, "bold": True, "color": PALETTE["text"], "space_after": 0}],
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.72),
            Inches(0.82),
            Inches(10.9),
            Inches(0.42),
            [{"text": subtitle, "size": 13, "color": PALETTE["muted"], "space_after": 0}],
        )

    add_pill(
        slide,
        Inches(11.48),
        Inches(0.36),
        Inches(1.02),
        Inches(0.34),
        f"{slide_no:02d}/{total:02d}",
        f"{ACCENTS[accent_idx % len(ACCENTS)]}_soft",
        font_size=12,
    )


def add_arrow(slide, left, top, width, height, fill_key: str = "blue_soft"):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
    set_shape_fill(arrow, fill_key)
    return arrow


def add_table(
    slide,
    left,
    top,
    width,
    height,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    col_widths: Sequence[float] | None = None,
) -> None:
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths is None:
        if len(headers) == 2:
            col_widths = [0.32, 0.68]
        elif len(headers) == 3:
            col_widths = [0.22, 0.33, 0.45]
        else:
            col_widths = [1 / len(headers)] * len(headers)
    for idx, ratio in enumerate(col_widths[: len(headers)]):
        table.columns[idx].width = int(width * ratio)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PALETTE["blue_soft"]
        cell.text_frame.paragraphs[0].font.name = FONT_HEAD
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.text_frame.paragraphs[0].font.color.rgb = PALETTE["text"]

    row_fill_cycle = ["surface", "surface_alt"]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = PALETTE[row_fill_cycle[(r - 1) % 2]]
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.font.color.rgb = PALETTE["text"]
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)


def build_cover(slide, total: int) -> None:
    add_full_background(slide, 0)
    add_pill(slide, Inches(0.75), Inches(0.48), Inches(1.6), Inches(0.36), "ODCR 汇报", "blue_soft", font_size=12)
    add_text(
        slide,
        Inches(0.75),
        Inches(1.0),
        Inches(7.8),
        Inches(1.8),
        [
            {"text": "第一次阶段汇报：", "size": 30, "bold": True, "space_after": 2},
            {"text": "证据锚定与内容 / 风格双通道预处理资产构建", "size": 26, "bold": True, "space_after": 10},
            {"text": "从 D4C 的原始文本预处理，推进到 ODCR 的显式证据锚定预处理", "size": 16, "color": PALETTE["muted"], "space_after": 0},
        ],
    )
    add_card(
        slide,
        Inches(0.75),
        Inches(3.1),
        Inches(7.6),
        Inches(1.55),
        "一句话摘要",
        [
            "本次汇报聚焦 ODCR 的预处理阶段，核心目标是把原始 review / explanation 文本",
            "升级为可供后续 shared / specific 解耦学习的证据资产。",
        ],
        BoxStyle(fill="surface", accent="mint", title_size=18, body_size=16),
    )
    add_card(
        slide,
        Inches(8.75),
        Inches(1.1),
        Inches(3.7),
        Inches(3.55),
        "汇报信息",
        [
            "课题：ODCR",
            "汇报人：zhangliml",
            "日期：2026-04-22",
            f"文稿来源：{SOURCE_TXT.name}",
        ],
        BoxStyle(fill="surface", accent="gold", title_size=18, body_size=16),
    )
    add_pill(slide, Inches(10.95), Inches(6.55), Inches(1.2), Inches(0.34), f"01/{total:02d}", "gold_soft", font_size=12)


def populate_overview_slide(slide, total: int, slide_no: int) -> None:
    add_slide_title(
        slide,
        "ODCR 整体创新链总览",
        "从 D4C 的 domain counterfactual augmentation，推进到“证据锚定—解耦表征—可靠性路由—受控生成”的完整链条",
        4,
        slide_no,
        total,
    )

    steps = [
        (
            "Preprocess",
            ["主创新：EASD 的证据锚点层", "辅创新：HSS 的风格先验", "作用：给样本显式打出内容证据与风格证据", "为后续 shared / specific 解耦提供输入来源"],
            "blue",
        ),
        (
            "Step3",
            ["主创新：EASD 的表征实现层", "辅创新：HSS 的层级风格结构", "作用：将内容与风格映射到 shared / specific 两个空间", "用结构化约束提升解耦效果"],
            "mint",
        ),
        (
            "Step4",
            ["主创新：RCR", "辅创新：UCI 元数据", "作用：评估反事实样本是否可靠", "决定样本送往评分路径还是解释路径"],
            "gold",
        ),
        (
            "Step5A",
            ["主创新：LCI", "作用：让评分主要依赖 shared 内容信息", "在反事实扰动下保持评分稳定", "降低风格噪声对推荐结果的干扰"],
            "coral",
        ),
        (
            "Step5B",
            ["主创新：CCV", "配套辅创新：UCI / FCA", "作用：区分“说什么”和“怎么说”", "让解释尽量与评分依赖的证据保持一致"],
            "lavender",
        ),
    ]

    left = Inches(0.55)
    gap = Inches(0.12)
    width = Inches(2.32)
    top = Inches(1.72)
    height = Inches(3.25)
    x = left
    for idx, (title, lines, accent) in enumerate(steps):
        add_card(
            slide,
            x,
            top,
            width,
            height,
            title,
            lines,
            BoxStyle(fill="surface", accent=accent, title_size=18, body_size=12),
        )
        if idx < len(steps) - 1:
            add_arrow(slide, x + width + Inches(0.02), Inches(3.02), Inches(0.28), Inches(0.22), "blue_soft")
        x += width + gap

    add_card(
        slide,
        Inches(0.72),
        Inches(5.45),
        Inches(12.0),
        Inches(1.0),
        "一句话总结",
        [
            "ODCR 的整体逻辑是：Preprocess 打锚点，Step3 做解耦，Step4 做可靠性路由，Step5A 保评分稳定，Step5B 做受控解释生成。",
            "因此第一次先讲 preprocess，因为后续所有模块都以前面的证据锚点为输入前提。",
        ],
        BoxStyle(fill="surface_alt", accent="blue", title_size=17, body_size=15),
    )


def build_slide_overview(slide, total: int) -> None:
    populate_overview_slide(slide, total, 2)


def build_slide_2(slide, total: int) -> None:
    add_slide_title(slide, "本次汇报回答 3 个问题", "从原始 D4C 到 ODCR 预处理资产化的最小闭环", 1, 3, total)
    specs = [
        (
            "Q1. 原始 D4C 的预处理怎么做？",
            ["review / explanation 提取", "5-core 过滤", "生成 processed.csv", "没有显式证据锚点层"],
            "blue",
        ),
        (
            "Q2. ODCR 的预处理改了什么？",
            ["新增内容锚点", "新增风格锚点", "新增 evidence quality 与 route 先验", "形成 content/style 双通道输入资产"],
            "mint",
        ),
        (
            "Q3. 这些资产后面怎么被用？",
            [
                "split / combine 保真传递",
                "compute_embeddings 构成 user / item 双通道画像",
                "infer_domain_semantics 构成 domain 双通道语义原型",
                "为 Step3 的 shared / specific 解耦做准备",
            ],
            "coral",
        ),
    ]
    left = Inches(0.72)
    gap = Inches(0.22)
    width = Inches(4.02)
    for idx, (title, lines, accent) in enumerate(specs):
        add_card(
            slide,
            left + idx * (width + gap),
            Inches(1.5),
            width,
            Inches(4.1),
            title,
            lines,
            BoxStyle(fill="surface", accent=accent, title_size=20, body_size=16),
        )
    add_card(
        slide,
        Inches(0.72),
        Inches(5.98),
        Inches(12.0),
        Inches(0.7),
        "收口",
        ["这一讲不是讲模型训练，而是回答“后续模型为什么能分得开”。"],
        BoxStyle(fill="surface_alt", accent="gold", title_size=16, body_size=16),
    )


def build_slide_3(slide, total: int) -> None:
    add_slide_title(slide, "原始 D4C 的预处理：原始文本进入后续模型", "左侧看流程，右侧看局限", 2, 4, total)
    add_card(
        slide,
        Inches(0.72),
        Inches(1.5),
        Inches(5.2),
        Inches(4.65),
        "原始流程",
        ["raw reviews", "抽取 review / explanation / rating / user / item", "5-core 过滤", "processed.csv"],
        BoxStyle(fill="surface", accent="blue", title_size=20, body_size=18),
    )
    for i in range(3):
        add_arrow(slide, Inches(2.87), Inches(2.3 + i * 0.92), Inches(0.36), Inches(0.26), "blue_soft")
    add_card(
        slide,
        Inches(6.18),
        Inches(1.5),
        Inches(6.55),
        Inches(4.65),
        "原始 D4C 的特点",
        [
            "预处理目标主要是“清洗与成表”",
            "没有显式 content evidence",
            "没有显式 style evidence",
            "没有 polarity anchor / local residual 等强语义资产",
            "本质还是“文本直接喂模型”",
        ],
        BoxStyle(fill="surface", accent="coral", title_size=20, body_size=17),
    )
    add_card(
        slide,
        Inches(0.72),
        Inches(6.05),
        Inches(12.0),
        Inches(0.72),
        "依据",
        ["原论文 D4C 的核心是 text attributes + domain counterfactual + factual/counterfactual joint training，而不是在预处理阶段显式构造 shared / specific 证据锚点。"],
        BoxStyle(fill="surface_alt", accent="mint", title_size=15, body_size=14),
    )


def build_slide_4(slide, total: int) -> None:
    add_slide_title(slide, "为什么原始预处理不够？", "把问题、后果和必要改动放在同一页说明白", 3, 5, total)
    rows = [
        (
            "问题 1：内容和风格混在同一段文本里",
            ["例子：房间很大，真的很舒服", "“大”更偏内容属性", "“真的很舒服”更偏表达风格 / 口吻"],
            ["后果：内容与风格信号纠缠", "后续 shared / specific 难以稳定分开"],
        ),
        (
            "问题 2：如果不提前打锚点",
            ["后续模型只能自己猜", "哪些部分是 item / content", "哪些部分是 style / domain expression"],
            ["后果：监督来源不明确", "解耦只能依赖隐式学习"],
        ),
        (
            "问题 3：只靠后续模型自动学",
            ["shared / specific 容易混淆", "风格可能污染评分", "后续解耦缺少证据来源"],
            ["后果：评分稳定性和解释可控性都会变差"],
        ),
    ]
    for idx, (title, left_lines, right_lines) in enumerate(rows):
        top = Inches(1.4 + idx * 1.45)
        add_card(
            slide,
            Inches(0.72),
            top,
            Inches(5.55),
            Inches(1.15),
            title,
            left_lines,
            BoxStyle(fill="surface", accent="gold", title_size=17, body_size=13),
        )
        add_arrow(slide, Inches(6.4), top + Inches(0.32), Inches(0.4), Inches(0.28), "gold_soft")
        add_card(
            slide,
            Inches(6.95),
            top,
            Inches(5.8),
            Inches(1.15),
            "后果",
            right_lines,
            BoxStyle(fill="surface", accent="coral", title_size=17, body_size=13),
        )
    add_card(
        slide,
        Inches(0.72),
        Inches(5.95),
        Inches(12.03),
        Inches(0.82),
        "结论",
        ["所以我在 preprocess 阶段做的事情，不是普通清洗，而是提前把“内容证据”和“风格证据”显式资产化。"],
        BoxStyle(fill="surface_alt", accent="mint", title_size=17, body_size=16),
    )


def build_slide_5(slide, total: int) -> None:
    add_slide_title(slide, "ODCR 预处理阶段：从“清洗数据”变成“构造证据资产”", "总图一页看完整链路", 4, 6, total)
    steps = [
        ("原始样本", ["user / item / rating", "review / explanation", "domain"], "blue"),
        ("证据锚定预处理", ["content anchor", "style anchor", "quality + route prior"], "mint"),
        (
            "新增锚点字段",
            [
                "content_keywords / aspects / entities / score",
                "style_markers / template / sentiment / domain_style_id / score",
                "evidence_quality / route_scorer / route_explainer",
            ],
            "gold",
        ),
        ("双通道资产构建", ["content channel", "style channel"], "coral"),
        ("供 Step3 使用", ["shared / specific disentanglement"], "lavender"),
    ]
    left = Inches(0.55)
    widths = [Inches(2.2), Inches(2.3), Inches(4.0), Inches(2.0), Inches(1.6)]
    x = left
    for idx, (title, lines, accent) in enumerate(steps):
        add_card(slide, x, Inches(2.0), widths[idx], Inches(3.1), title, lines, BoxStyle(fill="surface", accent=accent, title_size=19, body_size=14))
        x += widths[idx]
        if idx < len(steps) - 1:
            add_arrow(slide, x + Inches(0.05), Inches(3.08), Inches(0.35), Inches(0.28), "blue_soft")
            x += Inches(0.45)
    add_card(
        slide,
        Inches(0.75),
        Inches(5.55),
        Inches(11.9),
        Inches(0.92),
        "核心结论",
        ["根据审计，ODCR 的 preprocess 已经真实新增了一批显式锚点字段，这些字段会被 split / combine / embeddings / domain semantics / Step3 / Step5 继续消费。"],
        BoxStyle(fill="surface_alt", accent="blue", title_size=17, body_size=15),
    )


def build_slide_6(slide, total: int) -> None:
    add_slide_title(slide, "新增字段可以分为两类：内容锚点 与 风格锚点", "把字段、作用和说法边界同时讲清楚", 0, 7, total)
    rows = [
        ("内容锚点", "content_keywords", "抽取核心内容词"),
        ("内容锚点", "content_aspects", "提取 item aspect"),
        ("内容锚点", "content_entities", "保留实体信息"),
        ("内容锚点", "content_anchor_score", "形成内容强度分数"),
        ("风格锚点", "style_markers", "提取风格信号"),
        ("风格锚点", "template_family", "句式模板类别"),
        ("风格锚点", "sentiment_style", "评价语气 / 极性风格"),
        ("风格锚点", "length_style_bucket", "长度风格桶"),
        ("风格锚点", "domain_style_id", "域风格标识"),
        ("风格锚点", "style_anchor_score", "风格强度分数"),
        ("辅助字段", "evidence_quality", "证据质量"),
        ("辅助字段", "route_scorer / route_explainer", "路由先验"),
    ]
    add_table(slide, Inches(0.78), Inches(1.48), Inches(12.0), Inches(4.55), ["类别", "字段", "作用"], rows)
    add_card(
        slide,
        Inches(0.78),
        Inches(6.1),
        Inches(12.0),
        Inches(0.66),
        "表达边界",
        ["当前最稳的说法是“启发式证据锚定”，不要写成已经完成完整的强语义闭环。"],
        BoxStyle(fill="surface_alt", accent="coral", title_size=15, body_size=14),
    )


def build_slide_7(slide, total: int) -> None:
    add_slide_title(slide, "预处理阶段的形式化表达", "这页给老师一个“不是纯工程”的信号", 1, 8, total)
    cards = [
        (
            "1. 原始样本",
            ["x = (u, i, r, review, exp, d)", "u: user    i: item    r: rating", "exp: explanation    d: domain"],
            "blue",
        ),
        (
            "2. 内容 / 风格锚点映射",
            [
                "A_c(x) = ContentAnchor(x)",
                "A_s(x) = StyleAnchor(x)",
                "A_c(x) = {content_keywords, content_aspects, content_entities, content_anchor_score}",
                "A_s(x) = {style_markers, template_family, sentiment_style, domain_style_id, style_anchor_score}",
            ],
            "mint",
        ),
        (
            "3. 增强后的样本",
            ["x~ = (x, A_c(x), A_s(x), q(x), ρ(x))", "q(x): evidence quality", "ρ(x): route prior, 例如 scorer / explainer"],
            "gold",
        ),
    ]
    positions = [(Inches(0.78), Inches(1.58), Inches(3.7), Inches(3.75)), (Inches(4.8), Inches(1.58), Inches(4.0), Inches(3.75)), (Inches(9.1), Inches(1.58), Inches(3.7), Inches(3.75))]
    for (title, lines, accent), (left, top, width, height) in zip(cards, positions):
        add_card(slide, left, top, width, height, title, lines, BoxStyle(fill="surface", accent=accent, title_size=19, body_size=15))
    add_card(
        slide,
        Inches(0.78),
        Inches(5.75),
        Inches(12.0),
        Inches(0.82),
        "落点",
        ["也就是说，我不是只保留原样本，而是把样本扩展为“带证据锚点的增强样本”。"],
        BoxStyle(fill="surface_alt", accent="lavender", title_size=17, body_size=16),
    )


def build_slide_8(slide, total: int) -> None:
    add_slide_title(slide, "代码层面的真实落地证据", "这一页把“不是空想设计”讲明白", 2, 9, total)
    items = [
        (
            "1. 预处理落地",
            [
                "code/preprocess_data.py",
                "_build_odcr_anchors",
                "写出 content_keywords / content_aspects / content_entities / content_anchor_score",
                "写出 style_markers / template_family / sentiment_style / length_style_bucket",
                "写出 domain_style_id / style_anchor_score / evidence_quality / route_scorer / route_explainer",
            ],
            "blue",
        ),
        (
            "2. 字段保真传递",
            ["code/split_data.py", "code/combine_data.py", "作用不是创新本身", "但保证双域数据与锚点字段不会在中间丢失"],
            "mint",
        ),
        (
            "3. 下游真实消费",
            [
                "code/compute_embeddings.py",
                "生成 user_content_profiles.npy / user_style_profiles.npy",
                "生成 item_content_profiles.npy / item_style_profiles.npy",
                "code/infer_domain_semantics.py",
                "生成 domain_content.npy / domain_style.npy",
            ],
            "coral",
        ),
    ]
    for idx, (title, lines, accent) in enumerate(items):
        add_card(
            slide,
            Inches(0.72 + idx * 4.1),
            Inches(1.55),
            Inches(3.75),
            Inches(4.6),
            title,
            lines,
            BoxStyle(fill="surface", accent=accent, title_size=19, body_size=14),
        )
    add_card(
        slide,
        Inches(0.72),
        Inches(6.05),
        Inches(12.0),
        Inches(0.7),
        "结论",
        ["因此 preprocess 不是孤立字段堆砌，而是形成了下游真的会消费的双通道资产。"],
        BoxStyle(fill="surface_alt", accent="gold", title_size=16, body_size=15),
    )


def build_slide_9(slide, total: int) -> None:
    add_slide_title(slide, "实际例子：一条样本如何被预处理成“内容 / 风格双通道资产”", "这一页是最直观、最适合讲给老师看的例子", 3, 10, total)
    add_card(
        slide,
        Inches(0.72),
        Inches(1.35),
        Inches(12.0),
        Inches(0.9),
        "样本输入",
        ['Review / Explanation: "The rooms are spacious and the bathroom is a large tub."'],
        BoxStyle(fill="surface", accent="blue", title_size=17, body_size=16),
    )
    rows = [
        ("内容关键词", "rooms, bathroom, tub"),
        ("内容 aspect", "room size, bathroom facility"),
        ("内容实体", "room / bathroom"),
        ("内容分数", "0.82"),
        ("风格标记", "positive descriptive"),
        ("模板类型", "opinion + attribute"),
        ("情感风格", "positive"),
        ("长度桶", "short"),
        ("域风格 ID", "travel-domain"),
        ("风格分数", "0.67"),
        ("证据质量", "high"),
    ]
    add_table(slide, Inches(0.72), Inches(2.55), Inches(6.25), Inches(3.42), ["类型", "抽取结果"], rows)

    add_card(
        slide,
        Inches(7.2),
        Inches(2.55),
        Inches(5.52),
        Inches(1.15),
        "同一内容，不同风格",
        ['"The room is spacious."', '"Absolutely loved the roomy feel."', '"Rooms are quite large for the price."'],
        BoxStyle(fill="surface", accent="coral", title_size=17, body_size=15),
    )
    add_card(
        slide,
        Inches(7.2),
        Inches(3.95),
        Inches(5.52),
        Inches(0.95),
        "内容通道",
        ["回答“说了什么”"],
        BoxStyle(fill="surface_alt", accent="mint", title_size=17, body_size=16),
    )
    add_card(
        slide,
        Inches(7.2),
        Inches(5.02),
        Inches(5.52),
        Inches(0.95),
        "风格通道",
        ["回答“怎么说”"],
        BoxStyle(fill="surface_alt", accent="gold", title_size=17, body_size=16),
    )
    add_card(
        slide,
        Inches(0.72),
        Inches(6.1),
        Inches(12.0),
        Inches(0.6),
        "落点",
        ["只有先把内容与风格拆成双通道资产，后面 Step3 才有可能把 shared / specific 分开学。"],
        BoxStyle(fill="surface_alt", accent="lavender", title_size=15, body_size=14),
    )


def build_slide_10(slide, total: int) -> None:
    add_slide_title(slide, "预处理不是终点，而是 Step3 的输入前提", "这页把承接关系讲清楚", 4, 11, total)
    steps = [
        ("Preprocess", ["显式锚点字段构建"], "blue"),
        ("split / combine", ["字段保真"], "mint"),
        ("compute_embeddings", ["user_content_profiles.npy", "user_style_profiles.npy", "item_content_profiles.npy", "item_style_profiles.npy"], "gold"),
        ("infer_domain_semantics", ["domain_content.npy", "domain_style.npy"], "coral"),
        ("Step3", ["shared / specific disentanglement"], "lavender"),
    ]
    y = Inches(1.45)
    for idx, (title, lines, accent) in enumerate(steps):
        height = Inches(0.95 if idx != 2 else 1.2)
        add_card(slide, Inches(1.2), y, Inches(10.9), height, title, lines, BoxStyle(fill="surface", accent=accent, title_size=18, body_size=14))
        y += height + Inches(0.18)
        if idx < len(steps) - 1:
            add_arrow(slide, Inches(6.45), y - Inches(0.1), Inches(0.42), Inches(0.22), "blue_soft")
            y += Inches(0.14)
    add_card(
        slide,
        Inches(0.78),
        Inches(6.1),
        Inches(12.0),
        Inches(0.62),
        "关键信息",
        ["ODCR 预处理阶段的目标不是直接提升评分，而是为 Step3 的 shared 提供内容来源、为 specific 提供风格来源，并为后续解耦和路由建立数据基础。"],
        BoxStyle(fill="surface_alt", accent="mint", title_size=15, body_size=14),
    )


def build_slide_11(slide, total: int) -> None:
    add_slide_title(slide, "第一次汇报结论", "把这次工作的边界、价值和术语一次收住", 0, 12, total)
    conclusions = [
        ("结论 1", ["相较于原始 D4C 的“文本清洗式预处理”，ODCR 已把 preprocess 升级为证据资产构建阶段。"], "blue"),
        ("结论 2", ["当前代码里已经真实落地了显式锚点字段，以及 content/style 双通道的 user / item / domain 资产。"], "mint"),
        ("结论 3", ["目前最稳的术语不是“完整 EASD / HSS 已闭环”，而是“启发式证据锚定与内容 / 风格双通道预处理资产构建”。"], "gold"),
    ]
    for idx, (title, lines, accent) in enumerate(conclusions):
        add_card(
            slide,
            Inches(0.72 + idx * 4.06),
            Inches(1.72),
            Inches(3.82),
            Inches(3.95),
            title,
            lines,
            BoxStyle(fill="surface", accent=accent, title_size=21, body_size=17),
        )
    add_card(
        slide,
        Inches(0.72),
        Inches(5.95),
        Inches(12.0),
        Inches(0.72),
        "推荐收口",
        ["第一次工作的核心，不是训练更复杂的模型，而是把混杂文本整理成“内容可分、风格可分、域语义可分”的双通道证据资产。"],
        BoxStyle(fill="surface_alt", accent="coral", title_size=16, body_size=15),
    )


def build_slide_12(slide, total: int) -> None:
    add_slide_title(slide, "下一步：Step3 的共享 / 特异结构化解耦训练", "把下一次汇报的期待点提前埋好", 1, 13, total)
    add_card(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(7.15),
        Inches(4.45),
        "下一次继续回答",
        [
            "这些内容 / 风格双通道资产，如何进入 shared / specific 分支？",
            "Step3 的正交约束、shared invariance、specific separation 怎么实现？",
            "为什么它比原始 D4C 的单主干 + adv 对齐更适合解释任务？",
        ],
        BoxStyle(fill="surface", accent="blue", title_size=20, body_size=16),
    )
    add_card(
        slide,
        Inches(8.15),
        Inches(1.55),
        Inches(4.57),
        Inches(2.0),
        "建议关键词",
        ["Orthogonality", "Shared Invariance", "Specific Separation", "Structure-aware Disentanglement"],
        BoxStyle(fill="surface", accent="mint", title_size=18, body_size=15),
    )
    add_card(
        slide,
        Inches(8.15),
        Inches(3.78),
        Inches(4.57),
        Inches(2.22),
        "过渡句",
        [
            "这一轮先把“输入证据资产”站稳，下一轮再解释“模型为什么真的能分开”。",
            "这样逻辑最顺，也更符合当前代码与审计结论。",
        ],
        BoxStyle(fill="surface", accent="gold", title_size=18, body_size=15),
    )
    add_card(
        slide,
        Inches(0.72),
        Inches(6.16),
        Inches(12.0),
        Inches(0.56),
        "最终一句话",
        ["先把证据锚定做好，再去讲 shared / specific 解耦，汇报链条会更稳、更亮。"],
        BoxStyle(fill="surface_alt", accent="coral", title_size=15, body_size=14),
    )


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    prs.core_properties.title = "ODCR Preprocess Bright Report"
    prs.core_properties.subject = "ODCR preprocess presentation generated from TXT outline"
    prs.core_properties.author = "Codex"

    builders = [
        build_cover,
        build_slide_overview,
        build_slide_2,
        build_slide_3,
        build_slide_4,
        build_slide_5,
        build_slide_6,
        build_slide_7,
        build_slide_8,
        build_slide_9,
        build_slide_10,
        build_slide_11,
        build_slide_12,
    ]

    total = len(builders)
    blank = prs.slide_layouts[6]
    for builder in builders:
        slide = prs.slides.add_slide(blank)
        builder(slide, total)
    return prs


def build_single_slide_overview_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "ODCR Overall Innovation Chain"
    prs.core_properties.subject = "Single-slide overview extracted from TXT outline"
    prs.core_properties.author = "Codex"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    populate_overview_slide(slide, 1, 1)
    return prs


def main() -> None:
    if not SOURCE_TXT.exists():
        raise FileNotFoundError(f"未找到源文稿：{SOURCE_TXT}")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUT_PATH)
    single_prs = build_single_slide_overview_presentation()
    single_prs.save(SINGLE_SLIDE_OUT_PATH)
    print(f"saved: {OUT_PATH}")
    print(f"saved: {SINGLE_SLIDE_OUT_PATH}")


if __name__ == "__main__":
    main()
